import io
import math
import os
import zipfile
from dataclasses import dataclass
from typing import Dict, List, Tuple

import pandas as pd
from matplotlib import pyplot as plt
from matplotlib.gridspec import GridSpec
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


TITLE_AR = "قياس رضا أصحاب المصلحة عن برنامج هندسة الميكاترونيات"
RESPONSE_ORDER = [
    "موافق تماماً",
    "موافق",
    "محايد",
    "غير موافق",
    "غير موافق تماماً",
]
RESPONSE_COLORS = {
    "موافق تماماً": "#0B3C6D",     # Dark Blue
    "موافق": "#F28E2B",            # Orange
    "محايد": "#9AA0A6",            # Gray
    "غير موافق": "#F6C343",        # Yellow
    "غير موافق تماماً": "#63B3ED",  # Light Blue
}


def _shape_arabic(text: str) -> str:
    """Best-effort Arabic shaping without hard dependency."""
    if not isinstance(text, str):
        text = str(text)
    try:
        import arabic_reshaper  # type: ignore
        from bidi.algorithm import get_display  # type: ignore

        return get_display(arabic_reshaper.reshape(text))
    except Exception:
        return text


def _setup_matplotlib_fonts() -> None:
    plt.rcParams["font.family"] = [
        "Tahoma",
        "Arial",
        "Noto Naskh Arabic",
        "Amiri",
        "DejaVu Sans",
        "sans-serif",
    ]
    plt.rcParams["axes.unicode_minus"] = False


def load_data(file_bytes: bytes) -> pd.DataFrame:
    """Load survey workbook; first sheet is used."""
    return pd.read_excel(io.BytesIO(file_bytes), engine="openpyxl")


def analyze_question(series: pd.Series) -> Dict[str, float]:
    """Return response percentage distribution for one question."""
    clean = series.dropna().astype(str).str.strip()
    total = len(clean)
    if total == 0:
        return {category: 0.0 for category in RESPONSE_ORDER}

    counts = clean.value_counts()
    return {
        category: (float(counts.get(category, 0)) / total) * 100.0
        for category in RESPONSE_ORDER
    }


def calculate_satisfaction(df: pd.DataFrame) -> float:
    """Satisfaction = (agree + strongly agree) / total responses."""
    all_values = df.stack(dropna=True).astype(str).str.strip()
    total = len(all_values)
    if total == 0:
        return 0.0
    satisfied = int((all_values == "موافق").sum() + (all_values == "موافق تماماً").sum()) 
    return (satisfied / total) * 100.0


def _question_satisfaction(distribution: Dict[str, float]) -> float:
    return distribution.get("موافق", 0.0) + distribution.get("موافق تماماً", 0.0)


def _generate_insights(question_stats: Dict[str, Dict[str, float]], overall_satisfaction: float) -> Tuple[List[str], List[str]]:
    if not question_stats:
        return (["لا توجد بيانات كافية لاستخراج مؤشرات."], ["استكمال جمع الردود قبل التحليل."])

    ranked = sorted(
        ((q, _question_satisfaction(dist), dist) for q, dist in question_stats.items()),
        key=lambda x: x[1],
        reverse=True,
    )
    top_q, top_sat, _ = ranked[0]
    low_q, low_sat, low_dist = ranked[-1]
    avg_neutral = sum(d.get("محايد", 0.0) for _, _, d in ranked) / len(ranked)

    insights = [
        f"إجمالي الرضا العام: {overall_satisfaction:.1f}%",
        f"أعلى بند رضا: \"{top_q}\" ({top_sat:.1f}%).",
        f"أقل بند رضا: \"{low_q}\" ({low_sat:.1f}%).",
    ]
    if avg_neutral >= 20:
        insights.append(f"متوسط الحياد مرتفع نسبيًا ({avg_neutral:.1f}%).")

    improvement_points = []
    if low_sat < 70:
        improvement_points.append(f"التركيز على تحسين البند: \"{low_q}\".")
    if low_dist.get("غير موافق", 0.0) + low_dist.get("غير موافق تماماً", 0.0) > 20:
        improvement_points.append("تنفيذ إجراءات تصحيحية سريعة للفجوات الأكثر تكرارًا.")
    if not improvement_points:
        improvement_points.append("الاستمرار في تعزيز ممارسات التحسين المستمر مع متابعة دورية.")

    return insights[:4], improvement_points[:2]


def generate_charts(df: pd.DataFrame) -> Dict[str, Dict[str, float]]:
    """Build per-question distributions."""
    question_stats: Dict[str, Dict[str, float]] = {}
    for col in df.columns:
        question_stats[str(col)] = analyze_question(df[col])
    return question_stats


@dataclass
class SurveyResult:
    survey_name: str
    question_stats: Dict[str, Dict[str, float]]
    overall_satisfaction: float
    dashboard_png_path: str
    insights: List[str]
    improvements: List[str]


def build_dashboard(
    survey_name: str,
    question_stats: Dict[str, Dict[str, float]],
    overall_satisfaction: float,
    output_path: str,
) -> Tuple[List[str], List[str]]:
    """Create one high-resolution dashboard PNG per survey."""
    _setup_matplotlib_fonts()
    n_questions = max(1, len(question_stats))
    ncols_questions = 2
    nrows = math.ceil(n_questions / ncols_questions)

    fig = plt.figure(figsize=(22, 5.5 * nrows), dpi=180)
    gs = GridSpec(
        nrows=nrows,
        ncols=3,
        figure=fig,
        width_ratios=[1, 1, 0.85],
        wspace=0.28,
        hspace=0.48,
    )
    fig.patch.set_facecolor("#F8FAFC")

    title = f"{TITLE_AR}\n{survey_name}"
    fig.suptitle(_shape_arabic(title), fontsize=20, fontweight="bold", y=0.99)

    categories = RESPONSE_ORDER
    colors = [RESPONSE_COLORS[c] for c in categories]
    question_items = list(question_stats.items())

    for idx in range(nrows * ncols_questions):
        r, c = divmod(idx, ncols_questions)
        ax = fig.add_subplot(gs[r, c])
        if idx < len(question_items):
            question, distribution = question_items[idx]
            values = [distribution.get(cat, 0.0) for cat in categories]
            ax.bar(range(len(categories)), values, color=colors, width=0.68)
            ax.set_ylim(0, 100)
            ax.set_xticks(range(len(categories)))
            ax.set_xticklabels([_shape_arabic(c) for c in categories], rotation=25, ha="right", fontsize=9)
            ax.set_ylabel("Percentage %", fontsize=9)
            ax.set_title(_shape_arabic(question), fontsize=11, fontweight="semibold", pad=10)
            ax.grid(axis="y", linestyle="--", alpha=0.25)
            for i, val in enumerate(values):
                ax.text(i, val + 1.2, f"{val:.1f}%", ha="center", va="bottom", fontsize=8)
        else:
            ax.axis("off")

    insights, improvements = _generate_insights(question_stats, overall_satisfaction)

    summary_ax = fig.add_subplot(gs[:, 2])
    summary_ax.set_facecolor("#FFFFFF")
    for spine in summary_ax.spines.values():
        spine.set_visible(False)
    summary_ax.set_xticks([])
    summary_ax.set_yticks([])

    y = 0.95
    summary_ax.text(
        0.05,
        y,
        _shape_arabic("ملخص النتائج"),
        fontsize=16,
        fontweight="bold",
        va="top",
    )
    y -= 0.12
    summary_ax.text(
        0.05,
        y,
        _shape_arabic(f"نسبة الرضا الكلية: {overall_satisfaction:.1f}%"),
        fontsize=12.5,
        color="#0B3C6D",
        fontweight="bold",
        va="top",
    )

    y -= 0.10
    summary_ax.text(0.05, y, _shape_arabic("مؤشرات رئيسية:"), fontsize=12, fontweight="bold", va="top")
    y -= 0.06
    for insight in insights:
        summary_ax.text(0.08, y, _shape_arabic(f"• {insight}"), fontsize=10.5, va="top")
        y -= 0.055

    y -= 0.02
    summary_ax.text(0.05, y, _shape_arabic("نقاط تحسين:"), fontsize=12, fontweight="bold", va="top")
    y -= 0.06
    for point in improvements:
        summary_ax.text(0.08, y, _shape_arabic(f"• {point}"), fontsize=10.5, va="top")
        y -= 0.055

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    plt.tight_layout(rect=[0, 0, 1, 0.955])
    fig.savefig(output_path, dpi=300, bbox_inches="tight")
    plt.close(fig)
    return insights, improvements


def _set_arabic_font(run_or_font) -> None:
    """Configure font for proper Arabic/RTL text rendering."""
    run_or_font.name = "Tahoma"
    if hasattr(run_or_font, "latin"):
        run_or_font.latin = "Tahoma"
    if hasattr(run_or_font, "cs"):
        run_or_font.cs = "Tahoma"
    elif hasattr(run_or_font, "complex_script"):
        run_or_font.complex_script = "Tahoma"
    if hasattr(run_or_font, "_rPr"):
        rPr = run_or_font._rPr
        from lxml import etree
        cs_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}cs"
        cs_elem = rPr.find(cs_ns)
        if cs_elem is None:
            cs_elem = etree.SubElement(rPr, cs_ns)
        cs_elem.set("typeface", "Tahoma")
        if hasattr(run_or_font, "latin"):
            latin_elem = rPr.get_or_add_latin()
            latin_elem.set("typeface", "Tahoma")


def export_ppt(survey_results: List[SurveyResult], output_path: str) -> str:
    """Create PPT where each survey is a dedicated slide using native charts."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for result in survey_results:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.15), Inches(12.5), Inches(0.6))
        title_tf = title_box.text_frame
        title_tf.clear()
        run = title_tf.paragraphs[0].add_run()
        run.text = _shape_arabic(f"{TITLE_AR} - {result.survey_name}")
        run.font.size = Pt(19)
        run.font.bold = True
        _set_arabic_font(run.font)

        categories = RESPONSE_ORDER
        questions = list(result.question_stats.items())
        max_charts = min(6, len(questions))
        chart_positions = [
            (Inches(0.35), Inches(0.95)),
            (Inches(3.55), Inches(0.95)),
            (Inches(6.75), Inches(0.95)),
            (Inches(0.35), Inches(3.55)),
            (Inches(3.55), Inches(3.55)),
            (Inches(6.75), Inches(3.55)),
        ]
        chart_w = Inches(3.05)
        chart_h = Inches(2.25)

        for idx in range(max_charts):
            question, distribution = questions[idx]
            chart_data = CategoryChartData()
            chart_data.categories = [_shape_arabic(c) for c in categories]
            chart_data.add_series("النسبة", [distribution.get(cat, 0.0) for cat in categories])
            left, top = chart_positions[idx]
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                left,
                top,
                chart_w,
                chart_h,
                chart_data,
            )
            chart = chart_shape.chart
            chart.has_title = True
            chart.chart_title.text_frame.text = _shape_arabic(str(question)[:70])
            _set_arabic_font(chart.chart_title.text_frame.paragraphs[0].font)
            chart.value_axis.maximum_scale = 100
            chart.value_axis.minimum_scale = 0
            chart.value_axis.has_major_gridlines = True
            _set_arabic_font(chart.category_axis.tick_labels.font)
            _set_arabic_font(chart.value_axis.tick_labels.font)

            series = chart.series[0]
            for p_i, category in enumerate(categories):
                point = series.points[p_i]
                fill = point.format.fill
                fill.solid()
                hex_color = RESPONSE_COLORS[category].lstrip("#")
                fill.fore_color.rgb = RGBColor.from_string(hex_color)

        panel = slide.shapes.add_textbox(Inches(10.0), Inches(0.95), Inches(3.1), Inches(5.8))
        tf = panel.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = _shape_arabic(f"نسبة الرضا الكلية: {result.overall_satisfaction:.1f}%")
        p.font.bold = True
        p.font.size = Pt(14)
        _set_arabic_font(p.font)

        p2 = tf.add_paragraph()
        p2.text = _shape_arabic("مؤشرات رئيسية:")
        p2.font.bold = True
        p2.font.size = Pt(12)
        _set_arabic_font(p2.font)
        for insight in result.insights[:3]:
            pi = tf.add_paragraph()
            pi.text = _shape_arabic(f"• {insight}")
            pi.font.size = Pt(10.5)
            pi.level = 1
            _set_arabic_font(pi.font)

        p3 = tf.add_paragraph()
        p3.text = _shape_arabic("نقاط تحسين:")
        p3.font.bold = True
        p3.font.size = Pt(12)
        _set_arabic_font(p3.font)
        for item in result.improvements[:2]:
            pp = tf.add_paragraph()
            pp.text = _shape_arabic(f"• {item}")
            pp.font.size = Pt(10.5)
            pp.level = 1
            _set_arabic_font(pp.font)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path


def build_export_zip(files_to_add: List[str], output_zip_path: str) -> str:
    os.makedirs(os.path.dirname(output_zip_path), exist_ok=True)
    with zipfile.ZipFile(output_zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for path in files_to_add:
            zf.write(path, arcname=os.path.basename(path))
    return output_zip_path
